"""PPTX text extraction utilities."""

from pathlib import Path

from loguru import logger

from .text_extractor_common import TextExtractionError

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:  # pragma: no cover - optional dependency
    HAS_PPTX = False


def extract_text(file_path: Path) -> str:
    """Extract text from a PowerPoint presentation using python-pptx."""
    if not HAS_PPTX:
        raise TextExtractionError(
            "python-pptx is not available. Install it with: pip install python-pptx"
        )
    try:
        presentation = Presentation(str(file_path))
        extracted_text: list[str] = []
        if len(presentation.slides) == 0:
            logger.warning(f"PowerPoint presentation has no slides: {file_path}")
            return ""
        successful_slides = 0
        for slide_num, slide in enumerate(presentation.slides, 1):
            try:
                slide_content = [f"=== Slide {slide_num} ==="]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = _clean_slide_text(shape.text)
                        if text:
                            slide_content.append(text)
                if len(slide_content) > 1:
                    extracted_text.extend(slide_content)
                    extracted_text.append("")
                    successful_slides += 1
            except Exception as e:  # pragma: no cover - defensive
                logger.warning(
                    f"Failed to extract text from slide {slide_num} in {file_path}: {e}"
                )
                continue
        final_text = "\n".join(extracted_text).strip()
        if not final_text:
            if successful_slides == 0:
                raise TextExtractionError(
                    f"No text content could be extracted from PowerPoint: {file_path}. "
                    f"The presentation may contain only images or unsupported content."
                )
            logger.warning(
                f"No text content extracted from PowerPoint: {file_path}"
            )
            return ""
        logger.info(
            f"Successfully extracted text from PowerPoint: {file_path} "
            f"({successful_slides}/{len(presentation.slides)} slides processed)"
        )
        return final_text
    except FileNotFoundError as e:
        raise TextExtractionError(f"PowerPoint file not found: {file_path}") from e
    except PermissionError as e:
        raise TextExtractionError(
            f"Permission denied accessing PowerPoint file: {file_path}"
        ) from e
    except Exception as e:
        error_msg = str(e).lower()
        if "not a zip file" in error_msg or "bad zipfile" in error_msg:
            raise TextExtractionError(
                f"Invalid or corrupted PowerPoint file: {file_path}"
            ) from e
        if "password" in error_msg or "encrypted" in error_msg:
            raise TextExtractionError(
                f"Password-protected PowerPoint file: {file_path}"
            ) from e
        if "no such file" in error_msg or "package not found" in error_msg:
            raise TextExtractionError(f"PowerPoint file not found: {file_path}") from e
        raise TextExtractionError(
            f"Unexpected error extracting from PowerPoint {file_path}: {str(e)}"
        ) from e


def _clean_slide_text(text: str) -> str:
    """Clean and format slide text while preserving structure."""
    if not text or not text.strip():
        return ""
    lines = text.split("\n")
    cleaned_lines = []
    for line in lines:
        cleaned_line = line.strip()
        cleaned_lines.append(cleaned_line)
    import re

    cleaned_text = "\n".join(cleaned_lines)
    cleaned_text = re.sub(r"\n{3,}", "\n\n", cleaned_text)
    return cleaned_text.strip()
