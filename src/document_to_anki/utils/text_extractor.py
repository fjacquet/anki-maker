"""Text extraction utilities for various document formats."""

from pathlib import Path

from docx import Document
from loguru import logger
from pypdf import PdfReader

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pdfplumber

    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False


class TextExtractionError(Exception):
    """Exception raised when text extraction fails."""

    pass


class TextExtractor:
    """Handles text extraction from various document formats."""

    SUPPORTED_FORMATS = {".pdf", ".docx", ".pptx", ".txt", ".md"}

    def __init__(self) -> None:
        """Initialize the TextExtractor."""
        self.logger = logger

    def extract_text(self, file_path: Path) -> str:
        """
        Extract text from a file based on its extension.

        Args:
            file_path: Path to the file to extract text from

        Returns:
            Extracted text content as a string

        Raises:
            TextExtractionError: If extraction fails or format is unsupported
        """
        if not file_path.exists():
            raise TextExtractionError(f"File does not exist: {file_path}")

        if not file_path.is_file():
            raise TextExtractionError(f"Path is not a file: {file_path}")

        file_extension = file_path.suffix.lower()

        if file_extension not in self.SUPPORTED_FORMATS:
            raise TextExtractionError(
                f"Unsupported file format: {file_extension}. Supported formats: {', '.join(self.SUPPORTED_FORMATS)}"
            )

        try:
            if file_extension == ".pdf":
                return self.extract_text_from_pdf(file_path)
            elif file_extension == ".docx":
                return self.extract_text_from_docx(file_path)
            elif file_extension == ".pptx":
                return self.extract_text_from_pptx(file_path)
            elif file_extension in {".txt", ".md"}:
                return self.extract_text_from_txt(file_path)
            else:
                raise TextExtractionError(f"Unsupported file format: {file_extension}")

        except Exception as e:
            if isinstance(e, TextExtractionError):
                raise
            raise TextExtractionError(f"Failed to extract text from {file_path}: {str(e)}") from e

    def extract_text_from_pdf(self, file_path: Path) -> str:
        """
        Extract text from a PDF file using pypdf.

        Args:
            file_path: Path to the PDF file

        Returns:
            Extracted text content as a string

        Raises:
            TextExtractionError: If PDF extraction fails
        """
        try:
            text_content = []

            with open(file_path, "rb") as file:
                # Use strict=False to handle malformed PDFs more gracefully
                try:
                    pdf_reader = PdfReader(file)
                except Exception as reader_error:
                    # If PyPDF fails completely, try with different settings
                    self.logger.warning(f"Initial PDF reader failed for {file_path}: {reader_error}")
                    file.seek(0)  # Reset file pointer
                    try:
                        # Try with even more lenient settings
                        pdf_reader = PdfReader(file)
                    except Exception as fallback_error:
                        # If PyPDF fails completely, try pdfplumber as last resort
                        if HAS_PDFPLUMBER:
                            self.logger.warning(f"PyPDF failed for {file_path}, trying pdfplumber...")
                            return self._extract_with_pdfplumber(file_path)
                        else:
                            raise TextExtractionError(
                                f"Could not initialize PDF reader for {file_path}. "
                                f"The PDF may be severely corrupted or use an unsupported format. "
                                f"Original error: {reader_error}, Fallback error: {fallback_error}"
                            ) from fallback_error

                # Check if PDF is encrypted
                if pdf_reader.is_encrypted:
                    raise TextExtractionError(f"PDF file is encrypted: {file_path}")

                # Check if PDF has pages
                if len(pdf_reader.pages) == 0:
                    self.logger.warning(f"PDF file has no pages: {file_path}")
                    return ""

                # Extract text from each page with enhanced error handling
                successful_pages = 0
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():  # Only add non-empty pages
                            text_content.append(page_text)
                            successful_pages += 1
                    except (AttributeError, TypeError) as e:
                        # Handle NullObject and similar errors
                        self.logger.warning(f"Skipping malformed page {page_num + 1} in {file_path}: {e}")
                        continue
                    except Exception as e:
                        self.logger.warning(f"Failed to extract text from page {page_num + 1} in {file_path}: {e}")
                        continue

                extracted_text = "\n\n".join(text_content)

                if not extracted_text.strip():
                    if successful_pages == 0:
                        raise TextExtractionError(
                            f"No text content could be extracted from PDF: {file_path}. "
                            f"The PDF may be image-based, corrupted, or use unsupported formatting."
                        )
                    else:
                        self.logger.warning(f"No text content extracted from PDF: {file_path}")
                        return ""

                self.logger.info(
                    f"Successfully extracted text from PDF: {file_path} "
                    f"({successful_pages}/{len(pdf_reader.pages)} pages processed)"
                )
                return extracted_text

        except FileNotFoundError as e:
            raise TextExtractionError(f"PDF file not found: {file_path}") from e
        except PermissionError as e:
            raise TextExtractionError(f"Permission denied accessing PDF file: {file_path}") from e
        except Exception as e:
            # Handle pypdf specific errors
            if "PdfReadError" in str(type(e)) or "PDF" in str(e):
                raise TextExtractionError(f"Invalid or corrupted PDF file: {file_path} - {str(e)}") from e
            raise TextExtractionError(f"Unexpected error extracting from PDF {file_path}: {str(e)}") from e

    def _extract_with_pdfplumber(self, file_path: Path) -> str:
        """
        Extract text using pdfplumber as a fallback method.

        Args:
            file_path: Path to the PDF file

        Returns:
            Extracted text content as a string

        Raises:
            TextExtractionError: If pdfplumber extraction fails
        """
        if not HAS_PDFPLUMBER:
            raise TextExtractionError("pdfplumber is not available")

        try:
            text_content = []

            with pdfplumber.open(file_path) as pdf:
                if len(pdf.pages) == 0:
                    self.logger.warning(f"PDF file has no pages: {file_path}")
                    return ""

                successful_pages = 0
                for page_num, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_content.append(page_text)
                            successful_pages += 1
                    except Exception as e:
                        self.logger.warning(
                            f"pdfplumber failed to extract text from page {page_num + 1} in {file_path}: {e}"
                        )
                        continue

                extracted_text = "\n\n".join(text_content)

                if not extracted_text.strip():
                    if successful_pages == 0:
                        raise TextExtractionError(
                            f"No text content could be extracted from PDF using pdfplumber: {file_path}. "
                            f"The PDF may be image-based or use unsupported formatting."
                        )
                    else:
                        self.logger.warning(f"No text content extracted from PDF using pdfplumber: {file_path}")  # noqa: E501
                        return ""

                self.logger.info(
                    f"Successfully extracted text from PDF using pdfplumber: {file_path} "
                    f"({successful_pages}/{len(pdf.pages)} pages processed)"
                )
                return extracted_text

        except Exception as e:
            raise TextExtractionError(f"pdfplumber extraction failed for {file_path}: {str(e)}") from e

    def extract_text_from_docx(self, file_path: Path) -> str:
        """
        Extract text from a DOCX file using python-docx.

        Args:
            file_path: Path to the DOCX file

        Returns:
            Extracted text content as a string

        Raises:
            TextExtractionError: If DOCX extraction fails
        """
        try:
            doc = Document(str(file_path))
            text_content = []

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():  # Only add non-empty paragraphs
                    text_content.append(paragraph.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))

            extracted_text = "\n\n".join(text_content)

            if not extracted_text.strip():
                self.logger.warning(f"No text content extracted from DOCX: {file_path}")
                return ""

            self.logger.info(f"Successfully extracted text from DOCX: {file_path}")
            return extracted_text

        except FileNotFoundError as e:
            raise TextExtractionError(f"DOCX file not found: {file_path}") from e
        except PermissionError as e:
            raise TextExtractionError(f"Permission denied accessing DOCX file: {file_path}") from e
        except Exception as e:
            # Handle various docx-related errors
            error_msg = str(e).lower()
            if "not a zip file" in error_msg or "bad zipfile" in error_msg:
                raise TextExtractionError(f"Invalid or corrupted DOCX file: {file_path}") from e
            elif "no such file" in error_msg:
                raise TextExtractionError(f"DOCX file not found: {file_path}") from e
            else:
                raise TextExtractionError(f"Unexpected error extracting from DOCX {file_path}: {str(e)}") from e

    def extract_text_from_txt(self, file_path: Path) -> str:
        """
        Extract text from a TXT or MD file.

        Args:
            file_path: Path to the text file

        Returns:
            Extracted text content as a string

        Raises:
            TextExtractionError: If text extraction fails
        """
        try:
            # Try different encodings in order of preference
            encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]

            for encoding in encodings:
                try:
                    with open(file_path, encoding=encoding) as file:
                        content = file.read()

                    if not content.strip():
                        self.logger.warning(f"Text file is empty: {file_path}")
                        return ""

                    self.logger.info(
                        f"Successfully extracted text from {file_path.suffix.upper()} file: {file_path} "
                        f"(encoding: {encoding})"
                    )
                    return content

                except UnicodeDecodeError:
                    continue  # Try next encoding

            # If all encodings failed
            raise TextExtractionError(f"Could not decode text file with any supported encoding: {file_path}")

        except FileNotFoundError as e:
            raise TextExtractionError(f"Text file not found: {file_path}") from e
        except PermissionError as e:
            raise TextExtractionError(f"Permission denied accessing text file: {file_path}") from e
        except Exception as e:
            raise TextExtractionError(f"Unexpected error extracting from text file {file_path}: {str(e)}") from e

    def extract_text_from_pptx(self, file_path: Path) -> str:
        """
        Extract text from a PowerPoint presentation using python-pptx.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Extracted text content as a string with slide structure preserved

        Raises:
            TextExtractionError: If PPTX extraction fails
        """
        if not HAS_PPTX:
            raise TextExtractionError("python-pptx is not available. Install it with: pip install python-pptx")

        try:
            presentation = Presentation(str(file_path))
            extracted_text = []

            # Check if presentation has slides
            if len(presentation.slides) == 0:
                self.logger.warning(f"PowerPoint presentation has no slides: {file_path}")
                return ""

            successful_slides = 0
            for slide_num, slide in enumerate(presentation.slides, 1):
                try:
                    slide_content = [f"=== Slide {slide_num} ==="]

                    # Extract text from all shapes in the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            # Clean and preserve structure with bullet points and formatting
                            text = self._clean_slide_text(shape.text)
                            if text:
                                slide_content.append(text)

                    # Only add slides that have content beyond the header
                    if len(slide_content) > 1:
                        extracted_text.extend(slide_content)
                        extracted_text.append("")  # Add spacing between slides
                        successful_slides += 1

                except Exception as e:
                    self.logger.warning(f"Failed to extract text from slide {slide_num} in {file_path}: {e}")
                    continue

            final_text = "\n".join(extracted_text).strip()

            if not final_text:
                if successful_slides == 0:
                    raise TextExtractionError(
                        f"No text content could be extracted from PowerPoint: {file_path}. "
                        f"The presentation may contain only images or unsupported content."
                    )
                else:
                    self.logger.warning(f"No text content extracted from PowerPoint: {file_path}")
                    return ""

            self.logger.info(
                f"Successfully extracted text from PowerPoint: {file_path} "
                f"({successful_slides}/{len(presentation.slides)} slides processed)"
            )
            return final_text

        except FileNotFoundError as e:
            raise TextExtractionError(f"PowerPoint file not found: {file_path}") from e
        except PermissionError as e:
            raise TextExtractionError(f"Permission denied accessing PowerPoint file: {file_path}") from e
        except Exception as e:
            # Handle various pptx-related errors
            error_msg = str(e).lower()
            if "not a zip file" in error_msg or "bad zipfile" in error_msg:
                raise TextExtractionError(f"Invalid or corrupted PowerPoint file: {file_path}") from e
            elif "password" in error_msg or "encrypted" in error_msg:
                raise TextExtractionError(f"Password-protected PowerPoint file: {file_path}") from e
            elif "no such file" in error_msg or "package not found" in error_msg:
                raise TextExtractionError(f"PowerPoint file not found: {file_path}") from e
            else:
                raise TextExtractionError(f"Unexpected error extracting from PowerPoint {file_path}: {str(e)}") from e

    def _clean_slide_text(self, text: str) -> str:
        """
        Clean and format slide text while preserving structure.

        Args:
            text: Raw text from a slide shape

        Returns:
            Cleaned text with preserved structure

        """
        if not text or not text.strip():
            return ""

        # Split into lines and process each line
        lines = text.split("\n")
        cleaned_lines = []

        for line in lines:
            cleaned_line = line.strip()
            # Keep both empty and non-empty lines to preserve structure
            cleaned_lines.append(cleaned_line)

        # Join lines back together
        cleaned_text = "\n".join(cleaned_lines)

        # Remove multiple consecutive newlines but preserve intentional line breaks
        import re

        cleaned_text = re.sub(r"\n{3,}", "\n\n", cleaned_text)

        # Remove leading and trailing whitespace from the entire text
        cleaned_text = cleaned_text.strip()

        return cleaned_text

    def extract_text_from_md(self, file_path: Path) -> str:
        """
        Extract text from a Markdown file.

        This is an alias for extract_text_from_txt since both are plain text files.

        Args:
            file_path: Path to the markdown file

        Returns:
            Extracted text content as a string

        Raises:
            TextExtractionError: If text extraction fails
        """
        return self.extract_text_from_txt(file_path)

    def is_supported_format(self, file_path: Path) -> bool:
        """
        Check if a file format is supported for text extraction.

        Args:
            file_path: Path to the file to check

        Returns:
            True if the file format is supported, False otherwise
        """
        return file_path.suffix.lower() in self.SUPPORTED_FORMATS

    def get_supported_formats(self) -> set[str]:
        """
        Get the set of supported file formats.

        Returns:
            Set of supported file extensions (including the dot)
        """
        return self.SUPPORTED_FORMATS.copy()
